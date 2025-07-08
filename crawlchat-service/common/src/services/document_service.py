"""
Document service for handling document processing and management.
"""

import logging
import uuid
import asyncio
from datetime import datetime
from typing import Optional, List, Dict, Any
from pathlib import Path
import io
from functools import lru_cache
import hashlib

from common.src.models.documents import (
    Document, DocumentType, DocumentStatus, DocumentUpload, 
    DocumentProcessResponse, DocumentList
)
from common.src.core.exceptions import DocumentProcessingError
from common.src.services.unified_storage_service import unified_storage_service
from common.src.core.database import mongodb
from common.src.core.aws_config import aws_config

logger = logging.getLogger(__name__)

class DocumentService:
    """Document service for managing document processing and storage using MongoDB."""
    
    def __init__(self):
        # Cache for document type and content type mappings
        self._document_type_cache = {}
        self._content_type_cache = {}
    
    async def upload_document(self, upload_data: DocumentUpload) -> Document:
        """Upload a document with enhanced file integrity checking."""
        try:
            logger.info(f"📤 DOCUMENT UPLOAD START: {upload_data.filename}")
            
            # 🔍 FILE INTEGRITY CHECK - Log file details at upload
            file_content = upload_data.file_content
            logger.info(f"🔍 UPLOAD FILE INTEGRITY CHECK:")
            logger.info(f"   📏 File size: {len(file_content):,} bytes")
            logger.info(f"   🔢 File hash (MD5): {hashlib.md5(file_content).hexdigest()}")
            logger.info(f"   📄 First 20 bytes: {file_content[:20]}")
            logger.info(f"   📄 Last 20 bytes: {file_content[-20:]}")
            
            # Check if it's a valid PDF
            if upload_data.filename.lower().endswith('.pdf'):
                if file_content.startswith(b'%PDF-'):
                    logger.info(f"   ✅ Valid PDF header detected")
                else:
                    logger.warning(f"   ❌ Invalid PDF header: {file_content[:10]}")
                
                if b'%%EOF' in file_content[-1000:]:
                    logger.info(f"   ✅ PDF EOF marker found")
                else:
                    logger.warning(f"   ❌ PDF EOF marker missing")
            
            # Determine document type
            file_extension = Path(upload_data.filename).suffix.lower()
            document_type = self._get_document_type(file_extension)
            
            # Generate unique document ID
            document_id = str(uuid.uuid4())
            
            # Store file using unified storage service
            result = await unified_storage_service.upload_user_document(
                file_content=file_content,
                filename=upload_data.filename,
                user_id=upload_data.user_id,
                content_type=self._get_content_type(file_extension)
            )
            file_path = result["s3_key"]
            
            # Create document record
            document = Document(
                document_id=document_id,
                user_id=upload_data.user_id,
                filename=upload_data.filename,
                file_path=file_path,
                document_type=document_type,
                file_size=len(file_content),
                status=DocumentStatus.UPLOADED,
                uploaded_at=datetime.utcnow()
            )
            
            # Save to database
            await mongodb.connect()
            await mongodb.get_collection("documents").insert_one(document.dict())
            
            logger.info(f"✅ DOCUMENT UPLOAD COMPLETE: {document_id}")
            return document
            
        except Exception as e:
            logger.error(f"❌ Document upload failed: {e}")
            raise DocumentProcessingError(f"Failed to upload document: {e}")
    
    async def process_preprocessed_document(self, document_id: str, user_id: str, preprocessed_s3_key: str) -> DocumentProcessResponse:
        """Process a document that has been preprocessed by the preprocessing service."""
        try:
            await mongodb.connect()
            
            start_time = datetime.utcnow()
            
            # Get document data
            doc_data = await mongodb.get_collection("documents").find_one({"document_id": document_id, "user_id": user_id})
            if not doc_data:
                raise DocumentProcessingError("Document not found")

            document = Document(**doc_data)
            
            # Update status to processing
            await mongodb.get_collection("documents").update_one(
                {"document_id": document.document_id},
                {"$set": {"status": DocumentStatus.PROCESSING}}
            )

            # Extract content from preprocessed document using Textract
            content = await self._extract_preprocessed_content(preprocessed_s3_key, document.filename)
            
            # Generate summary and key points in parallel
            summary_task = self._generate_summary(content, None)
            key_points_task = self._extract_key_points(content)
            
            summary, key_points = await asyncio.gather(summary_task, key_points_task)

            # Update document with all changes at once
            update_data = {
                "content": content,
                "summary": summary,
                "key_points": key_points,
                "status": DocumentStatus.PROCESSED,
                "processed_at": datetime.utcnow(),
                "preprocessed_s3_key": preprocessed_s3_key
            }
            
            await mongodb.get_collection("documents").update_one(
                {"document_id": document.document_id},
                {"$set": update_data}
            )
            
            processing_time = (datetime.utcnow() - start_time).total_seconds()
            logger.info(f"Processed preprocessed document {document.document_id} in {processing_time:.2f}s")
            
            return DocumentProcessResponse(
                document_id=document.document_id,
                status=DocumentStatus.PROCESSED,
                content=content,
                summary=summary,
                key_points=key_points,
                processing_time=processing_time
            )
            
        except Exception as e:
            logger.error(f"Error processing preprocessed document: {e}")
            raise DocumentProcessingError(f"Failed to process preprocessed document: {e}")

    async def _extract_preprocessed_content(self, s3_key: str, filename: str) -> str:
        """Extract content from a preprocessed document using Textract."""
        try:
            from common.src.services.aws_textract_service import textract_service, DocumentType
            
            logger.info(f"Extracting content from preprocessed document: {s3_key}")
            
            # Determine document type
            document_type = DocumentType.GENERAL
            filename_lower = filename.lower()
            if any(keyword in filename_lower for keyword in ['form', 'invoice', 'receipt', 'tax', 'w2', '1099']):
                document_type = DocumentType.FORM
                logger.info(f"Detected form/invoice document: {filename}")
            
            # Process with Textract
            bucket_name = aws_config.s3_bucket
            text_content, page_count = await textract_service.process_preprocessed_document(
                bucket_name, s3_key, document_type
            )
            
            # Store page count for later use
            self._last_page_count = page_count
            
            if text_content and text_content.strip():
                logger.info(f"Successfully extracted {len(text_content)} characters from preprocessed document: {filename}")
                return text_content
            else:
                logger.warning(f"No content extracted from preprocessed document: {filename}")
                return f"Could not extract text from preprocessed document: {filename}"
                
        except Exception as e:
            logger.error(f"Error extracting content from preprocessed document {s3_key}: {e}")
            return f"Error extracting content from preprocessed document: {e}"

    async def process_document(self, document_path: str, user_query: Optional[str] = None) -> DocumentProcessResponse:
        """Process a document and extract content with optimized performance."""
        try:
            await mongodb.connect()
            
            start_time = datetime.utcnow()
            
            # Get document data
            doc_data = await mongodb.get_collection("documents").find_one({"file_path": document_path})
            if not doc_data:
                raise DocumentProcessingError("Document not found")

            document = Document(**doc_data)
            
            # Early return if already processed
            if document.status == DocumentStatus.PROCESSED and document.content:
                logger.info(f"Document {document.document_id} already processed, returning cached content")
                return DocumentProcessResponse(
                    document_id=document.document_id,
                    status=DocumentStatus.PROCESSED,
                    content=document.content,
                    summary=document.summary or "",
                    key_points=document.key_points or [],
                    processing_time=0.0
                )
            
            # Update status to processing
            await mongodb.get_collection("documents").update_one(
                {"document_id": document.document_id},
                {"$set": {"status": DocumentStatus.PROCESSING}}
            )

            # Extract content
            self._last_page_count = None
            content = await self._extract_content(document)
            
            # Check page limit before storing the document
            page_count = getattr(self, '_last_page_count', None)
            if page_count is not None:
                # Check if this would exceed the user's page limit
                if not await self.check_page_limit_for_user(document.user_id, page_count):
                    # Revert status back to uploaded
                    await mongodb.get_collection("documents").update_one(
                        {"document_id": document.document_id},
                        {"$set": {"status": DocumentStatus.UPLOADED}}
                    )
                    current_total = await self.get_total_pages_for_user(document.user_id)
                    raise DocumentProcessingError(
                        f"You have already processed {current_total} pages. "
                        f"This document has {page_count} pages, which would exceed your 100-page limit. "
                        f"Please upgrade your plan to process more pages."
                    )
                
                # Store page count in metadata
                document.metadata['page_count'] = page_count
                await mongodb.get_collection("documents").update_one(
                    {"document_id": document.document_id},
                    {"$set": {"metadata.page_count": page_count}}
                )
            
            # Generate summary and key points in parallel
            summary_task = self._generate_summary(content, user_query)
            key_points_task = self._extract_key_points(content)
            
            summary, key_points = await asyncio.gather(summary_task, key_points_task)

            # Update document with all changes at once
            update_data = {
                "content": content,
                "summary": summary,
                "key_points": key_points,
                "status": DocumentStatus.PROCESSED,
                "processed_at": datetime.utcnow()
            }
            
            await mongodb.get_collection("documents").update_one(
                {"document_id": document.document_id},
                {"$set": update_data}
            )
            
            processing_time = (datetime.utcnow() - start_time).total_seconds()
            logger.info(f"Processed document {document.document_id} in {processing_time:.2f}s")
            
            return DocumentProcessResponse(
                document_id=document.document_id,
                status=DocumentStatus.PROCESSED,
                content=content,
                summary=summary,
                key_points=key_points,
                processing_time=processing_time
            )
            
        except Exception as e:
            logger.error(f"Error processing document: {e}")
            raise DocumentProcessingError(f"Failed to process document: {e}")
    
    async def _extract_content(self, document: Document) -> str:
        """Extract content from document based on type with optimized performance."""
        try:
            storage_service = get_storage_service()
            
            content = await storage_service.get_file_content(document.file_path)
            if not content:
                logger.error(f"No content found for file: {document.filename}")
                return f"Could not read file: {document.filename}"
            
            logger.info(f"Processing document: {document.filename} (type: {document.document_type})")
            
            # Fast path for text-based documents
            if document.document_type in {DocumentType.TXT, DocumentType.HTML, DocumentType.JSON}:
                try:
                    text_content = content.decode('utf-8', errors='ignore')
                    logger.info(f"Successfully extracted text from {document.filename}")
                    return text_content
                except Exception as e:
                    logger.error(f"Error decoding text content from {document.filename}: {e}")
                    return f"Error decoding text content: {e}"
            
            # CSV processing
            if document.document_type == DocumentType.CSV:
                try:
                    import csv
                    text_content = content.decode('utf-8', errors='ignore')
                    # Convert CSV to readable text format
                    csv_reader = csv.reader(text_content.splitlines())
                    rows = list(csv_reader)
                    if rows:
                        # Convert to tabular text format
                        formatted_rows = []
                        for i, row in enumerate(rows):
                            formatted_rows.append(f"Row {i+1}: {' | '.join(str(cell) for cell in row)}")
                        result = '\n'.join(formatted_rows)
                        logger.info(f"Successfully extracted CSV content from {document.filename}")
                        return result
                    else:
                        return "Empty CSV file"
                except Exception as e:
                    logger.error(f"Error processing CSV {document.filename}: {e}")
                    return f"Error processing CSV: {e}"
            
            # PDF processing with improved error handling
            if document.document_type == DocumentType.PDF:
                return await self._extract_pdf_text_optimized(content, document.filename)
            
            # Excel files (XLSX, XLS)
            if document.document_type in {DocumentType.XLSX, DocumentType.XLS}:
                try:
                    from openpyxl import load_workbook
                    # Load workbook from bytes
                    wb = load_workbook(io.BytesIO(content), data_only=True)
                    text_content = []
                    
                    for sheet_name in wb.sheetnames:
                        ws = wb[sheet_name]
                        text_content.append(f"Sheet: {sheet_name}")
                        
                        for row in ws.iter_rows(values_only=True):
                            if any(cell is not None for cell in row):
                                row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                                text_content.append(row_text)
                    
                    result = '\n'.join(text_content)
                    logger.info(f"Successfully extracted Excel content from {document.filename}")
                    return result
                except Exception as e:
                    logger.error(f"Error processing Excel file {document.filename}: {e}")
                    return f"Error processing Excel file: {e}"
            
            # PowerPoint files (PPT, PPTX)
            if document.document_type in {DocumentType.PPT, DocumentType.PPTX}:
                try:
                    from pptx import Presentation
                    prs = Presentation(io.BytesIO(content))
                    text_content = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text_content.append(shape.text)
                    result = '\n'.join(text_content)
                    if result.strip():
                        logger.info(f"Successfully extracted PowerPoint content from {document.filename}")
                        return result
                    else:
                        return "No text content found in PowerPoint file"
                except Exception as e:
                    logger.error(f"Error processing PowerPoint file {document.filename}: {e}")
                    return f"Error processing PowerPoint file: {e}"
            
            # Word documents (DOC, DOCX)
            if document.document_type in {DocumentType.DOC, DocumentType.DOCX}:
                try:
                    import docx
                    doc = docx.Document(io.BytesIO(content))
                    text_content = []
                    for paragraph in doc.paragraphs:
                        text_content.append(paragraph.text)
                    result = '\n'.join(text_content)
                    if result.strip():
                        logger.info(f"Successfully extracted Word document content from {document.filename}")
                        return result
                    else:
                        return "No text content found in Word document"
                except Exception as e:
                    logger.error(f"Error processing Word document {document.filename}: {e}")
                    return f"Error processing Word document: {e}"
            
            # Fallback for unknown types
            logger.warning(f"Unknown document type: {document.document_type} for {document.filename}")
            return f"Content extraction not implemented for {document.document_type}"
                
        except Exception as e:
            logger.error(f"Error extracting content from {document.filename}: {e}")
            return f"Error extracting content: {e}"
    
    async def _extract_pdf_text_optimized(self, content: bytes, filename: str) -> str:
        """Optimized PDF text extraction with early returns and OCR fallback."""
        
        logger.info(f"Starting PDF extraction for: {filename}")
        
        # Try PyPDF2 first (fastest)
        try:
            import PyPDF2
            pdf = PyPDF2.PdfReader(io.BytesIO(content))
            text_content = []
            total_pages = len(pdf.pages)
            
            for i, page in enumerate(pdf.pages):
                try:
                    page_text = page.extract_text()
                    if page_text and page_text.strip():
                        text_content.append(f"Page {i+1}: {page_text}")
                    else:
                        logger.warning(f"Page {i+1} appears to be empty or image-based")
                except Exception as e:
                    logger.warning(f"Failed to extract text from page {i+1}: {e}")
            
            if text_content:
                result = '\n\n'.join(text_content)
                logger.info(f"Successfully extracted PDF text using PyPDF2: {filename} ({total_pages} pages, {len(text_content)} pages with text)")
                return result
            else:
                logger.warning(f"PyPDF2 extracted no text from {filename} - file may be image-based or corrupted")
        except Exception as e:
            logger.warning(f"PyPDF2 failed for {filename}: {e}")
        
        # Try pdfminer.six (more robust)
        try:
            from pdfminer.high_level import extract_text
            text_content = extract_text(io.BytesIO(content))
            if text_content and text_content.strip():
                logger.info(f"Successfully extracted PDF text using pdfminer.six: {filename}")
                return text_content
            else:
                logger.warning(f"pdfminer.six extracted no text from {filename} - file may be image-based or corrupted")
        except Exception as e:
            logger.warning(f"pdfminer.six failed for {filename}: {e}")
        
        # Try AWS Textract for scanned PDFs
        try:
            textract_text = await self._extract_text_with_ocr(content, filename)
            if textract_text and textract_text.strip():
                logger.info(f"Successfully extracted text using AWS Textract: {filename}")
                return textract_text
            else:
                logger.warning(f"AWS Textract extracted no text from {filename}")
        except Exception as e:
            logger.warning(f"AWS Textract failed for {filename}: {e}")
        
        # Final fallback - try to decode as text
        try:
            text_content = content.decode('utf-8', errors='ignore')
            if text_content.strip():
                logger.info(f"Fallback text extraction successful for {filename}")
                return text_content
        except Exception as e:
            logger.error(f"Final fallback failed for {filename}: {e}")
        
        logger.error(f"All PDF extraction methods failed for {filename}")
        return f"Could not extract text from PDF: {filename}. This may be because:\n- The PDF is image-based (scanned document)\n- The PDF is corrupted or damaged\n- The PDF has no embedded text content\n- AWS Textract is not available or configured\n\nPlease try uploading a text-based PDF or a different document format."
    
    async def _extract_text_with_ocr(self, content: bytes, filename: str) -> str:
        """Extract text from PDF using AWS Textract for scanned documents."""
        try:
            from common.src.services.aws_textract_service import textract_service, DocumentType
            
            logger.info(f"Starting AWS Textract extraction for: {filename}")
            
            # Determine document type (default to general)
            document_type = DocumentType.GENERAL
            
            # Simple heuristic to detect form/invoice documents
            # In a production system, you might want more sophisticated detection
            filename_lower = filename.lower()
            if any(keyword in filename_lower for keyword in ['form', 'invoice', 'receipt', 'tax', 'w2', '1099']):
                document_type = DocumentType.FORM
                logger.info(f"Detected form/invoice document: {filename}")
            
            # Use the new Textract service
            text_content, page_count = await textract_service.upload_to_s3_and_extract(
                content, 
                filename, 
                document_type,
                "document_service"  # Use a generic user_id for document service calls
            )
            
            self._last_page_count = page_count  # Store for later use
            if text_content and text_content.strip():
                logger.info(f"Successfully extracted text using AWS Textract: {filename}")
                return text_content
            else:
                logger.warning(f"AWS Textract extracted no text from {filename}")
                return ""
                
        except ImportError as e:
            logger.warning(f"Textract service not available: {e}")
            return "Textract not available: Required dependencies are not installed. Please install AWS SDK."
        except Exception as e:
            logger.error(f"Textract extraction failed for {filename}: {e}")
            return f"Textract extraction failed: {e}"

    async def _generate_summary(self, content: str, user_query: Optional[str]) -> str:
        """Generate summary of document content."""
        if user_query:
            return f"Summary based on query '{user_query}': {content[:200]}..."
            return f"Document summary: {content[:200]}..."
    
    async def _extract_key_points(self, content: str) -> List[str]:
        """Extract key points from document content."""
        return [
            "Key point 1 (placeholder)",
            "Key point 2 (placeholder)",
            "Key point 3 (placeholder)"
        ]
    
    def _get_document_type(self, extension: str) -> DocumentType:
        """Get document type from file extension with caching."""
        if extension not in self._document_type_cache:
            extension_map = {
                '.pdf': DocumentType.PDF,
                '.doc': DocumentType.DOC,
                '.docx': DocumentType.DOCX,
                '.txt': DocumentType.TXT,
                '.html': DocumentType.HTML,
                '.htm': DocumentType.HTML,
                '.xlsx': DocumentType.XLSX,
                '.xls': DocumentType.XLS,
                '.csv': DocumentType.CSV,
                '.ppt': DocumentType.PPT,
                '.pptx': DocumentType.PPTX,
                '.json': DocumentType.JSON,
                '.xml': DocumentType.TXT,
                '.md': DocumentType.TXT,
                '.rtf': DocumentType.TXT,
            }
            self._document_type_cache[extension] = extension_map.get(extension, DocumentType.TXT)
        return self._document_type_cache[extension]
    
    def _get_content_type(self, extension: str) -> str:
        """Get content type from file extension with caching."""
        if extension not in self._content_type_cache:
            content_types = {
                '.pdf': 'application/pdf',
                '.doc': 'application/msword',
                '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                '.txt': 'text/plain',
                '.html': 'text/html',
                '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                '.xls': 'application/vnd.ms-excel',
                '.csv': 'text/csv',
                '.ppt': 'application/vnd.ms-powerpoint',
                '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                '.json': 'application/json',
            }
            self._content_type_cache[extension] = content_types.get(extension, 'application/octet-stream')
        return self._content_type_cache[extension]
    
    async def get_document(self, document_id: str, user_id: str) -> Optional[Document]:
        """Get a document by ID."""
        try:
            await mongodb.connect()
            doc = await mongodb.get_collection("documents").find_one({"document_id": document_id, "user_id": user_id})
            return Document(**doc) if doc else None
        except Exception as e:
            logger.error(f"Error getting document {document_id}: {e}")
            return None
    
    async def list_user_documents(self, user_id: str, limit: int = 50, skip: int = 0) -> DocumentList:
        """List documents for a user with optimized querying."""
        try:
            await mongodb.connect()
            # Use projection to only fetch needed fields
            projection = {
                "document_id": 1,
                "user_id": 1,
                "filename": 1,
                "file_path": 1,
                "document_type": 1,
                "status": 1,
                "uploaded_at": 1,
                "processed_at": 1,
                "file_size": 1
            }
            cursor = mongodb.get_collection("documents").find(
                {"user_id": user_id}, 
                projection=projection
            ).skip(skip).limit(limit)
            docs = [Document(**doc) async for doc in cursor]
            total_count = await mongodb.get_collection("documents").count_documents({"user_id": user_id})
            return DocumentList(documents=docs, total_count=total_count)
        except Exception as e:
            logger.error(f"Error listing documents for user {user_id}: {e}")
            return DocumentList(documents=[], total_count=0)
    
    async def delete_document(self, document_id: str, user_id: str) -> bool:
        """Delete a document with optimized operations."""
        try:
            await mongodb.connect()
            
            doc_data = await mongodb.get_collection("documents").find_one(
                {"document_id": document_id, "user_id": user_id},
                projection={"file_path": 1}
            )
            if not doc_data:
                return False
            
            # Parallel delete operations
            delete_tasks = []
            
            # Delete from storage
            try:
                storage_service = get_storage_service()
                delete_tasks.append(storage_service.delete_file(doc_data["file_path"]))
                logger.info(f"Deleted file: {doc_data['file_path']}")
            except Exception as e:
                logger.warning(f"Failed to delete file {doc_data['file_path']}: {e}")
            
            # Delete from database
            delete_tasks.append(
                mongodb.get_collection("documents").delete_one({"document_id": document_id, "user_id": user_id})
            )
            
            # Wait for all delete operations
            results = await asyncio.gather(*delete_tasks, return_exceptions=True)
            
            # Check database deletion result
            db_result = results[-1]
            if isinstance(db_result, Exception):
                logger.error(f"Database deletion failed: {db_result}")
                return False
            
            return db_result.deleted_count > 0
        except Exception as e:
            logger.error(f"Error deleting document {document_id}: {e}")
            return False
    
    async def delete_document_by_filename(self, filename: str, user_id: str) -> bool:
        """Delete a document by filename with optimized operations."""
        try:
            await mongodb.connect()
            
            doc_data = await mongodb.get_collection("documents").find_one(
                {"filename": filename, "user_id": user_id},
                projection={"document_id": 1, "file_path": 1}
            )
            if not doc_data:
                return False
            
            # Parallel delete operations
            delete_tasks = []
            
            # Delete from storage
            try:
                storage_service = get_storage_service()
                delete_tasks.append(storage_service.delete_file(doc_data["file_path"]))
                logger.info(f"Deleted file: {doc_data['file_path']}")
            except Exception as e:
                logger.warning(f"Failed to delete file {doc_data['file_path']}: {e}")
            
            # Delete from database
            delete_tasks.append(
                mongodb.get_collection("documents").delete_one({"document_id": doc_data["document_id"], "user_id": user_id})
            )
            
            # Wait for all delete operations
            results = await asyncio.gather(*delete_tasks, return_exceptions=True)
            
            # Check database deletion result
            db_result = results[-1]
            if isinstance(db_result, Exception):
                logger.error(f"Database deletion failed: {db_result}")
                return False
            
            return db_result.deleted_count > 0
        except Exception as e:
            logger.error(f"Error deleting document by filename {filename}: {e}")
            return False
    
    async def get_task_documents(self, task_id: str) -> List[Document]:
        """Get documents for a specific task with optimized querying."""
        try:
            # Use projection to only fetch needed fields
            projection = {
                "document_id": 1,
                "user_id": 1,
                "filename": 1,
                "file_path": 1,
                "file_size": 1,
                "document_type": 1,
                "status": 1,
                "uploaded_at": 1,
                "task_id": 1
            }
            
            # First try to get documents from the documents collection
            cursor = mongodb.get_collection("documents").find(
                {"task_id": task_id}, 
                projection=projection
            )
            documents = [Document(**doc) async for doc in cursor]
            
            if documents:
                return documents
            
            # If no documents found, try crawl_tasks collection
            crawl_task = await mongodb.get_collection("crawl_tasks").find_one(
                {"task_id": task_id},
                projection={"user_id": 1, "documents": 1}
            )
            
            if crawl_task and crawl_task.get("documents"):
                # Convert crawl task documents to Document objects
                documents = []
                for doc_data in crawl_task["documents"]:
                    document = Document(
                        document_id=doc_data.get("document_id", str(uuid.uuid4())),
                        user_id=crawl_task.get("user_id", ""),
                        filename=doc_data.get("filename", ""),
                        file_path=doc_data.get("file_path", ""),
                        file_size=doc_data.get("file_size", 0),
                        document_type=self._get_document_type(Path(doc_data.get("filename", "")).suffix),
                        status=DocumentStatus.PROCESSED,
                        uploaded_at=doc_data.get("uploaded_at", datetime.utcnow()),
                        task_id=task_id,
                        crawl_task_id=task_id
                    )
                    documents.append(document)
                
                return documents
            
            return []
            
        except Exception as e:
            logger.error(f"Error getting task documents for {task_id}: {e}")
            return []
    
    async def create_document(self, document: Document) -> Document:
        """Create a new document record in the database."""
        try:
            await mongodb.get_collection("documents").insert_one(document.dict())
            logger.info(f"Created document record: {document.document_id}")
            return document
        except Exception as e:
            logger.error(f"Error creating document record: {e}")
            raise DocumentProcessingError(f"Failed to create document record: {e}")

    async def get_total_pages_for_user(self, user_id: str) -> int:
        """Get total number of pages processed by a user across all documents."""
        try:
            await mongodb.connect()
            # Sum up page_count from all processed documents for this user
            pipeline = [
                {"$match": {"user_id": user_id, "status": "processed"}},
                {"$group": {"_id": None, "total_pages": {"$sum": "$metadata.page_count"}}}
            ]
            cursor = mongodb.get_collection("documents").aggregate(pipeline)
            result = await cursor.to_list(length=None)
            total_pages = result[0]["total_pages"] if result else 0
            logger.info(f"User {user_id} has processed {total_pages} total pages")
            return total_pages
        except Exception as e:
            logger.error(f"Error getting total pages for user {user_id}: {e}")
            return 0

    async def check_page_limit_for_user(self, user_id: str, new_page_count: int) -> bool:
        """Check if processing a document with new_page_count would exceed the 100-page limit."""
        try:
            current_total = await self.get_total_pages_for_user(user_id)
            new_total = current_total + new_page_count
            if new_total > 100:
                logger.warning(f"User {user_id} would exceed 100-page limit: {current_total} + {new_page_count} = {new_total}")
                return False
            return True
        except Exception as e:
            logger.error(f"Error checking page limit for user {user_id}: {e}")
            return False

    async def process_preprocessed_document(self, document_id: str, user_id: str, preprocessed_s3_key: str) -> DocumentProcessResponse:
        """Process a document that has been preprocessed by the preprocessing service."""
        try:
            await mongodb.connect()
            
            start_time = datetime.utcnow()
            
            # Get document data
            doc_data = await mongodb.get_collection("documents").find_one({"document_id": document_id, "user_id": user_id})
            if not doc_data:
                raise DocumentProcessingError("Document not found")

            document = Document(**doc_data)
            
            # Update status to processing
            await mongodb.get_collection("documents").update_one(
                {"document_id": document.document_id},
                {"$set": {"status": DocumentStatus.PROCESSING}}
            )

            # Extract content from preprocessed document using Textract
            content = await self._extract_preprocessed_content(preprocessed_s3_key, document.filename)
            
            # Generate summary and key points in parallel
            summary_task = self._generate_summary(content, None)
            key_points_task = self._extract_key_points(content)
            
            summary, key_points = await asyncio.gather(summary_task, key_points_task)

            # Update document with all changes at once
            update_data = {
                "content": content,
                "summary": summary,
                "key_points": key_points,
                "status": DocumentStatus.PROCESSED,
                "processed_at": datetime.utcnow(),
                "preprocessed_s3_key": preprocessed_s3_key
            }
            
            await mongodb.get_collection("documents").update_one(
                {"document_id": document.document_id},
                {"$set": update_data}
            )
            
            processing_time = (datetime.utcnow() - start_time).total_seconds()
            logger.info(f"Processed preprocessed document {document.document_id} in {processing_time:.2f}s")
            
            return DocumentProcessResponse(
                document_id=document.document_id,
                status=DocumentStatus.PROCESSED,
                content=content,
                summary=summary,
                key_points=key_points,
                processing_time=processing_time
            )
            
        except Exception as e:
            logger.error(f"Error processing preprocessed document: {e}")
            raise DocumentProcessingError(f"Failed to process preprocessed document: {e}")

    async def _extract_preprocessed_content(self, s3_key: str, filename: str) -> str:
        """Extract content from a preprocessed document using Textract."""
        try:
            from common.src.services.aws_textract_service import textract_service, DocumentType
            
            logger.info(f"Extracting content from preprocessed document: {s3_key}")
            
            # Determine document type
            document_type = DocumentType.GENERAL
            filename_lower = filename.lower()
            if any(keyword in filename_lower for keyword in ['form', 'invoice', 'receipt', 'tax', 'w2', '1099']):
                document_type = DocumentType.FORM
                logger.info(f"Detected form/invoice document: {filename}")
            
            # Process with Textract
            bucket_name = aws_config.s3_bucket
            text_content, page_count = await textract_service.process_preprocessed_document(
                bucket_name, s3_key, document_type
            )
            
            # Store page count for later use
            self._last_page_count = page_count
            
            if text_content and text_content.strip():
                logger.info(f"Successfully extracted {len(text_content)} characters from preprocessed document: {filename}")
                return text_content
            else:
                logger.warning(f"No content extracted from preprocessed document: {filename}")
                return f"Could not extract text from preprocessed document: {filename}"
                
        except Exception as e:
            logger.error(f"Error extracting content from preprocessed document {s3_key}: {e}")
            return f"Error extracting content from preprocessed document: {e}"

# Global instance
document_service = DocumentService() 