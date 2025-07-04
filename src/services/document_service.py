"""
Document service for handling document processing and management.
"""

import logging
import uuid
import asyncio
from datetime import datetime
from typing import Optional, List, Dict, Any
from pathlib import Path
import io
from functools import lru_cache

from src.models.documents import (
    Document, DocumentType, DocumentStatus, DocumentUpload, 
    DocumentProcessResponse, DocumentList
)
from src.core.exceptions import DocumentProcessingError
from src.services.storage_service import get_storage_service
from src.core.database import mongodb

logger = logging.getLogger(__name__)

class DocumentService:
    """Document service for managing document processing and storage using MongoDB."""
    
    def __init__(self):
        # Cache for document type and content type mappings
        self._document_type_cache = {}
        self._content_type_cache = {}
    
    async def upload_document(self, upload_data: DocumentUpload) -> Document:
        """Upload and store a document with optimized performance."""
        try:
            await mongodb.connect()
            
            document_id = str(uuid.uuid4())
            extension = Path(upload_data.filename).suffix.lower()
            document_type = self._get_document_type(extension)
            storage_service = get_storage_service()
            
            # Generate S3 key more efficiently
            timestamp = int(datetime.utcnow().timestamp())
            unique_id = str(uuid.uuid4())[:8]
            s3_key = f"uploaded_documents/{upload_data.user_id}/{timestamp}_{unique_id}{extension}"

            # Create document record first to avoid waiting for S3 upload
            document = Document(
                document_id=document_id,
                user_id=upload_data.user_id,
                filename=upload_data.filename,
                file_path=s3_key,  # Use S3 key directly
                file_size=len(upload_data.content),  # Use actual size
                document_type=document_type,
                status=DocumentStatus.UPLOADED,
                uploaded_at=datetime.utcnow(),
                task_id=upload_data.task_id
            )
            
            # Parallel operations: S3 upload and database insert
            upload_task = storage_service.upload_file_from_bytes(
                upload_data.content, s3_key, content_type=self._get_content_type(extension)
            )
            db_task = mongodb.get_collection("documents").insert_one(document.dict())
            
            # Wait for both operations to complete
            await asyncio.gather(upload_task, db_task)
            
            logger.info(f"Uploaded document {document_id}: {upload_data.filename}")
            return document
            
        except Exception as e:
            logger.error(f"Error uploading document: {e}")
            raise DocumentProcessingError(f"Failed to upload document: {e}")
    
    async def process_document(self, document_path: str, user_query: Optional[str] = None) -> DocumentProcessResponse:
        """Process a document and extract content with optimized performance."""
        try:
            await mongodb.connect()
            
            start_time = datetime.utcnow()
            
            # Get document data
            doc_data = await mongodb.get_collection("documents").find_one({"file_path": document_path})
            if not doc_data:
                raise DocumentProcessingError("Document not found")

            document = Document(**doc_data)
            
            # Early return if already processed
            if document.status == DocumentStatus.PROCESSED and document.content:
                logger.info(f"Document {document.document_id} already processed, returning cached content")
                return DocumentProcessResponse(
                    document_id=document.document_id,
                    status=DocumentStatus.PROCESSED,
                    content=document.content,
                    summary=document.summary or "",
                    key_points=document.key_points or [],
                    processing_time=0.0
                )
            
            # Update status to processing
            await mongodb.get_collection("documents").update_one(
                {"document_id": document.document_id},
                {"$set": {"status": DocumentStatus.PROCESSING}}
            )

            # Extract content
            content = await self._extract_content(document)
            
            # Generate summary and key points in parallel
            summary_task = self._generate_summary(content, user_query)
            key_points_task = self._extract_key_points(content)
            
            summary, key_points = await asyncio.gather(summary_task, key_points_task)

            # Update document with all changes at once
            update_data = {
                "content": content,
                "summary": summary,
                "key_points": key_points,
                "status": DocumentStatus.PROCESSED,
                "processed_at": datetime.utcnow()
            }
            
            await mongodb.get_collection("documents").update_one(
                {"document_id": document.document_id},
                {"$set": update_data}
            )
            
            processing_time = (datetime.utcnow() - start_time).total_seconds()
            logger.info(f"Processed document {document.document_id} in {processing_time:.2f}s")
            
            return DocumentProcessResponse(
                document_id=document.document_id,
                status=DocumentStatus.PROCESSED,
                content=content,
                summary=summary,
                key_points=key_points,
                processing_time=processing_time
            )
            
        except Exception as e:
            logger.error(f"Error processing document: {e}")
            raise DocumentProcessingError(f"Failed to process document: {e}")
    
    async def _extract_content(self, document: Document) -> str:
        """Extract content from document based on type with optimized performance."""
        try:
            storage_service = get_storage_service()
            
            content = await storage_service.get_file_content(document.file_path)
            if not content:
                logger.error(f"No content found for file: {document.filename}")
                return f"Could not read file: {document.filename}"
            
            logger.info(f"Processing document: {document.filename} (type: {document.document_type})")
            
            # Fast path for text-based documents
            if document.document_type in {DocumentType.TXT, DocumentType.HTML, DocumentType.JSON}:
                try:
                    text_content = content.decode('utf-8', errors='ignore')
                    logger.info(f"Successfully extracted text from {document.filename}")
                    return text_content
                except Exception as e:
                    logger.error(f"Error decoding text content from {document.filename}: {e}")
                    return f"Error decoding text content: {e}"
            
            # CSV processing
            if document.document_type == DocumentType.CSV:
                try:
                    import csv
                    text_content = content.decode('utf-8', errors='ignore')
                    # Convert CSV to readable text format
                    csv_reader = csv.reader(text_content.splitlines())
                    rows = list(csv_reader)
                    if rows:
                        # Convert to tabular text format
                        formatted_rows = []
                        for i, row in enumerate(rows):
                            formatted_rows.append(f"Row {i+1}: {' | '.join(str(cell) for cell in row)}")
                        result = '\n'.join(formatted_rows)
                        logger.info(f"Successfully extracted CSV content from {document.filename}")
                        return result
                    else:
                        return "Empty CSV file"
                except Exception as e:
                    logger.error(f"Error processing CSV {document.filename}: {e}")
                    return f"Error processing CSV: {e}"
            
            # PDF processing with improved error handling
            if document.document_type == DocumentType.PDF:
                return await self._extract_pdf_text_optimized(content, document.filename)
            
            # Excel files (XLSX, XLS)
            if document.document_type in {DocumentType.XLSX, DocumentType.XLS}:
                try:
                    from openpyxl import load_workbook
                    # Load workbook from bytes
                    wb = load_workbook(io.BytesIO(content), data_only=True)
                    text_content = []
                    
                    for sheet_name in wb.sheetnames:
                        ws = wb[sheet_name]
                        text_content.append(f"Sheet: {sheet_name}")
                        
                        for row in ws.iter_rows(values_only=True):
                            if any(cell is not None for cell in row):
                                row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                                text_content.append(row_text)
                    
                    result = '\n'.join(text_content)
                    logger.info(f"Successfully extracted Excel content from {document.filename}")
                    return result
                except Exception as e:
                    logger.error(f"Error processing Excel file {document.filename}: {e}")
                    return f"Error processing Excel file: {e}"
            
            # PowerPoint files (PPT, PPTX)
            if document.document_type in {DocumentType.PPT, DocumentType.PPTX}:
                try:
                    from pptx import Presentation
                    prs = Presentation(io.BytesIO(content))
                    text_content = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text_content.append(shape.text)
                    result = '\n'.join(text_content)
                    if result.strip():
                        logger.info(f"Successfully extracted PowerPoint content from {document.filename}")
                        return result
                    else:
                        return "No text content found in PowerPoint file"
                except Exception as e:
                    logger.error(f"Error processing PowerPoint file {document.filename}: {e}")
                    return f"Error processing PowerPoint file: {e}"
            
            # Word documents (DOC, DOCX)
            if document.document_type in {DocumentType.DOC, DocumentType.DOCX}:
                try:
                    import docx
                    doc = docx.Document(io.BytesIO(content))
                    text_content = []
                    for paragraph in doc.paragraphs:
                        text_content.append(paragraph.text)
                    result = '\n'.join(text_content)
                    if result.strip():
                        logger.info(f"Successfully extracted Word document content from {document.filename}")
                        return result
                    else:
                        return "No text content found in Word document"
                except Exception as e:
                    logger.error(f"Error processing Word document {document.filename}: {e}")
                    return f"Error processing Word document: {e}"
            
            # Fallback for unknown types
            logger.warning(f"Unknown document type: {document.document_type} for {document.filename}")
            return f"Content extraction not implemented for {document.document_type}"
                
        except Exception as e:
            logger.error(f"Error extracting content from {document.filename}: {e}")
            return f"Error extracting content: {e}"
    
    async def _extract_pdf_text_optimized(self, content: bytes, filename: str) -> str:
        """Optimized PDF text extraction with early returns and OCR fallback."""
        
        logger.info(f"Starting PDF extraction for: {filename}")
        
        # Try PyPDF2 first (fastest)
        try:
            import PyPDF2
            pdf = PyPDF2.PdfReader(io.BytesIO(content))
            text_content = []
            total_pages = len(pdf.pages)
            
            for i, page in enumerate(pdf.pages):
                try:
                    page_text = page.extract_text()
                    if page_text and page_text.strip():
                        text_content.append(f"Page {i+1}: {page_text}")
                    else:
                        logger.warning(f"Page {i+1} appears to be empty or image-based")
                except Exception as e:
                    logger.warning(f"Failed to extract text from page {i+1}: {e}")
            
            if text_content:
                result = '\n\n'.join(text_content)
                logger.info(f"Successfully extracted PDF text using PyPDF2: {filename} ({total_pages} pages, {len(text_content)} pages with text)")
                return result
            else:
                logger.warning(f"PyPDF2 extracted no text from {filename} - file may be image-based or corrupted")
        except Exception as e:
            logger.warning(f"PyPDF2 failed for {filename}: {e}")
        
        # Try pdfminer.six (more robust)
        try:
            from pdfminer.high_level import extract_text
            text_content = extract_text(io.BytesIO(content))
            if text_content and text_content.strip():
                logger.info(f"Successfully extracted PDF text using pdfminer.six: {filename}")
                return text_content
            else:
                logger.warning(f"pdfminer.six extracted no text from {filename} - file may be image-based or corrupted")
        except Exception as e:
            logger.warning(f"pdfminer.six failed for {filename}: {e}")
        
        # Try OCR for scanned PDFs
        try:
            ocr_text = await self._extract_text_with_ocr(content, filename)
            if ocr_text and ocr_text.strip():
                logger.info(f"Successfully extracted text using OCR: {filename}")
                return ocr_text
            else:
                logger.warning(f"OCR extracted no text from {filename}")
        except Exception as e:
            logger.warning(f"OCR failed for {filename}: {e}")
        
        # Final fallback - try to decode as text
        try:
            text_content = content.decode('utf-8', errors='ignore')
            if text_content.strip():
                logger.info(f"Fallback text extraction successful for {filename}")
                return text_content
        except Exception as e:
            logger.error(f"Final fallback failed for {filename}: {e}")
        
        logger.error(f"All PDF extraction methods failed for {filename}")
        return f"Could not extract text from PDF: {filename}. This may be because:\n- The PDF is image-based (scanned document)\n- The PDF is corrupted or damaged\n- The PDF has no embedded text content\n\nPlease try uploading a text-based PDF or a different document format."
    
    async def _extract_text_with_ocr(self, content: bytes, filename: str) -> str:
        """Extract text from PDF using OCR for scanned documents."""
        try:
            import pytesseract
            from pdf2image import convert_from_bytes
            from PIL import Image
            import tempfile
            import os
            
            logger.info(f"Starting OCR extraction for: {filename}")
            
            # Check if Tesseract is available
            try:
                # Try to get Tesseract version to check if it's installed
                tesseract_version = pytesseract.get_tesseract_version()
                logger.info(f"Tesseract version: {tesseract_version}")
            except Exception as e:
                logger.error(f"Tesseract not available: {e}")
                # Try to find tesseract binary manually
                import subprocess
                try:
                    result = subprocess.run(['which', 'tesseract'], capture_output=True, text=True)
                    if result.returncode == 0:
                        logger.info(f"Tesseract found at: {result.stdout.strip()}")
                        # Try to set the path manually
                        pytesseract.pytesseract.tesseract_cmd = result.stdout.strip()
                        tesseract_version = pytesseract.get_tesseract_version()
                        logger.info(f"Tesseract version after manual path: {tesseract_version}")
                    else:
                        logger.error("Tesseract binary not found in PATH")
                        return f"Could not extract text from PDF: {filename}. This appears to be an image-based (scanned) document, but OCR is not available in this environment. Please try uploading a text-based PDF or contact support for OCR capabilities."
                except Exception as sub_e:
                    logger.error(f"Error checking tesseract binary: {sub_e}")
                    return f"Could not extract text from PDF: {filename}. This appears to be an image-based (scanned) document, but OCR is not available in this environment. Please try uploading a text-based PDF or contact support for OCR capabilities."
            
            # Check if Poppler is available by testing pdf2image
            try:
                # This will fail if poppler-utils is not installed
                from pdf2image.exceptions import PDFPageCountError
                logger.info("Poppler utilities are available")
                
                # Also check if pdftoppm is available
                import subprocess
                result = subprocess.run(['which', 'pdftoppm'], capture_output=True, text=True)
                if result.returncode == 0:
                    logger.info(f"pdftoppm found at: {result.stdout.strip()}")
                else:
                    logger.error("pdftoppm binary not found in PATH")
                    return "OCR is not available. Please contact support or try again later."
            except ImportError:
                logger.error("Poppler utilities not available")
                return "OCR is not available. Please contact support or try again later."
            
            # Convert PDF pages to images
            try:
                images = convert_from_bytes(content, dpi=300)  # Higher DPI for better OCR accuracy
                logger.info(f"Successfully converted PDF to {len(images)} images")
            except Exception as e:
                logger.error(f"Failed to convert PDF to images: {e}")
                return f"OCR failed: Could not convert PDF to images. Error: {e}"
            
            if not images:
                logger.warning(f"No images extracted from PDF: {filename}")
                return ""
            
            text_content = []
            
            for i, image in enumerate(images):
                try:
                    # Preprocess image for better OCR
                    # Convert to grayscale
                    gray_image = image.convert('L')
                    
                    # Use pytesseract to extract text
                    page_text = pytesseract.image_to_string(gray_image, lang='eng')
                    
                    if page_text and page_text.strip():
                        text_content.append(f"Page {i+1} (OCR): {page_text}")
                        logger.info(f"OCR successful for page {i+1}")
                    else:
                        logger.warning(f"OCR extracted no text from page {i+1}")
                        
                except Exception as e:
                    logger.warning(f"OCR failed for page {i+1}: {e}")
            
            if text_content:
                result = '\n\n'.join(text_content)
                logger.info(f"OCR extraction completed for {filename} ({len(images)} pages)")
                return result
            else:
                logger.warning(f"OCR extracted no text from any page in {filename}")
                return ""
                
        except ImportError as e:
            logger.warning(f"OCR dependencies not available: {e}")
            return "OCR not available: Required dependencies (pytesseract, pdf2image, PIL) are not installed. Please use a container-based Lambda function with OCR support."
        except Exception as e:
            logger.error(f"OCR extraction failed for {filename}: {e}")
            return f"OCR extraction failed: {e}"

    async def _generate_summary(self, content: str, user_query: Optional[str]) -> str:
        """Generate summary of document content."""
        if user_query:
            return f"Summary based on query '{user_query}': {content[:200]}..."
            return f"Document summary: {content[:200]}..."
    
    async def _extract_key_points(self, content: str) -> List[str]:
        """Extract key points from document content."""
        return [
            "Key point 1 (placeholder)",
            "Key point 2 (placeholder)",
            "Key point 3 (placeholder)"
        ]
    
    def _get_document_type(self, extension: str) -> DocumentType:
        """Get document type from file extension with caching."""
        if extension not in self._document_type_cache:
            extension_map = {
                '.pdf': DocumentType.PDF,
                '.doc': DocumentType.DOC,
                '.docx': DocumentType.DOCX,
                '.txt': DocumentType.TXT,
                '.html': DocumentType.HTML,
                '.htm': DocumentType.HTML,
                '.xlsx': DocumentType.XLSX,
                '.xls': DocumentType.XLS,
                '.csv': DocumentType.CSV,
                '.ppt': DocumentType.PPT,
                '.pptx': DocumentType.PPTX,
                '.json': DocumentType.JSON,
                '.xml': DocumentType.TXT,
                '.md': DocumentType.TXT,
                '.rtf': DocumentType.TXT,
            }
            self._document_type_cache[extension] = extension_map.get(extension, DocumentType.TXT)
        return self._document_type_cache[extension]
    
    def _get_content_type(self, extension: str) -> str:
        """Get content type from file extension with caching."""
        if extension not in self._content_type_cache:
            content_types = {
                '.pdf': 'application/pdf',
                '.doc': 'application/msword',
                '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                '.txt': 'text/plain',
                '.html': 'text/html',
                '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                '.xls': 'application/vnd.ms-excel',
                '.csv': 'text/csv',
                '.ppt': 'application/vnd.ms-powerpoint',
                '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                '.json': 'application/json',
            }
            self._content_type_cache[extension] = content_types.get(extension, 'application/octet-stream')
        return self._content_type_cache[extension]
    
    async def get_document(self, document_id: str, user_id: str) -> Optional[Document]:
        """Get a document by ID."""
        try:
            await mongodb.connect()
            doc = await mongodb.get_collection("documents").find_one({"document_id": document_id, "user_id": user_id})
            return Document(**doc) if doc else None
        except Exception as e:
            logger.error(f"Error getting document {document_id}: {e}")
            return None

    async def list_user_documents(self, user_id: str, limit: int = 50, skip: int = 0) -> DocumentList:
        """List documents for a user with optimized querying."""
        try:
            await mongodb.connect()
            
            # Use projection to only fetch needed fields
            projection = {
                "document_id": 1,
                "filename": 1,
                "document_type": 1,
                "status": 1,
                "uploaded_at": 1,
                "processed_at": 1,
                "file_size": 1
            }
            
            cursor = mongodb.get_collection("documents").find(
                {"user_id": user_id}, 
                projection=projection
            ).skip(skip).limit(limit)
            
            docs = [Document(**doc) async for doc in cursor]
            total_count = await mongodb.get_collection("documents").count_documents({"user_id": user_id})
            return DocumentList(documents=docs, total_count=total_count)
        except Exception as e:
            logger.error(f"Error listing documents for user {user_id}: {e}")
            return DocumentList(documents=[], total_count=0)
    
    async def delete_document(self, document_id: str, user_id: str) -> bool:
        """Delete a document with optimized operations."""
        try:
            await mongodb.connect()
            
            doc_data = await mongodb.get_collection("documents").find_one(
                {"document_id": document_id, "user_id": user_id},
                projection={"file_path": 1}
            )
            if not doc_data:
                return False
                
            # Parallel delete operations
            delete_tasks = []
                
            # Delete from storage
            try:
                storage_service = get_storage_service()
                delete_tasks.append(storage_service.delete_file(doc_data["file_path"]))
                logger.info(f"Deleted file: {doc_data['file_path']}")
            except Exception as e:
                logger.warning(f"Failed to delete file {doc_data['file_path']}: {e}")
                
            # Delete from database
            delete_tasks.append(
                mongodb.get_collection("documents").delete_one({"document_id": document_id, "user_id": user_id})
            )
            
            # Wait for all delete operations
            results = await asyncio.gather(*delete_tasks, return_exceptions=True)
            
            # Check database deletion result
            db_result = results[-1]
            if isinstance(db_result, Exception):
                logger.error(f"Database deletion failed: {db_result}")
                return False
            
            return db_result.deleted_count > 0
        except Exception as e:
            logger.error(f"Error deleting document {document_id}: {e}")
            return False
    
    async def delete_document_by_filename(self, filename: str, user_id: str) -> bool:
        """Delete a document by filename with optimized operations."""
        try:
            await mongodb.connect()
            
            doc_data = await mongodb.get_collection("documents").find_one(
                {"filename": filename, "user_id": user_id},
                projection={"document_id": 1, "file_path": 1}
            )
            if not doc_data:
                return False
                
            # Parallel delete operations
            delete_tasks = []
                
            # Delete from storage
            try:
                storage_service = get_storage_service()
                delete_tasks.append(storage_service.delete_file(doc_data["file_path"]))
                logger.info(f"Deleted file: {doc_data['file_path']}")
            except Exception as e:
                logger.warning(f"Failed to delete file {doc_data['file_path']}: {e}")
                
            # Delete from database
            delete_tasks.append(
                mongodb.get_collection("documents").delete_one({"document_id": doc_data["document_id"], "user_id": user_id})
            )
            
            # Wait for all delete operations
            results = await asyncio.gather(*delete_tasks, return_exceptions=True)
            
            # Check database deletion result
            db_result = results[-1]
            if isinstance(db_result, Exception):
                logger.error(f"Database deletion failed: {db_result}")
                return False
            
            return db_result.deleted_count > 0
        except Exception as e:
            logger.error(f"Error deleting document by filename {filename}: {e}")
            return False
    
    async def get_task_documents(self, task_id: str) -> List[Document]:
        """Get documents for a specific task with optimized querying."""
        try:
            # Use projection to only fetch needed fields
            projection = {
                "document_id": 1,
                "user_id": 1,
                "filename": 1,
                "file_path": 1,
                "file_size": 1,
                "document_type": 1,
                "status": 1,
                "uploaded_at": 1,
                "task_id": 1
            }
            
            # First try to get documents from the documents collection
            cursor = mongodb.get_collection("documents").find(
                {"task_id": task_id}, 
                projection=projection
            )
            documents = [Document(**doc) async for doc in cursor]
            
            if documents:
                return documents
            
            # If no documents found, try crawl_tasks collection
            crawl_task = await mongodb.get_collection("crawl_tasks").find_one(
                {"task_id": task_id},
                projection={"user_id": 1, "documents": 1}
            )
            
            if crawl_task and crawl_task.get("documents"):
                # Convert crawl task documents to Document objects
                documents = []
                for doc_data in crawl_task["documents"]:
                    document = Document(
                        document_id=doc_data.get("document_id", str(uuid.uuid4())),
                        user_id=crawl_task.get("user_id", ""),
                        filename=doc_data.get("filename", ""),
                        file_path=doc_data.get("file_path", ""),
                        file_size=doc_data.get("file_size", 0),
                        document_type=self._get_document_type(Path(doc_data.get("filename", "")).suffix),
                        status=DocumentStatus.PROCESSED,
                        uploaded_at=doc_data.get("uploaded_at", datetime.utcnow()),
                        task_id=task_id,
                        crawl_task_id=task_id
                    )
                    documents.append(document)
                
                return documents
            
            return []
            
        except Exception as e:
            logger.error(f"Error getting task documents for {task_id}: {e}")
            return []
    
    async def create_document(self, document: Document) -> Document:
        """Create a new document record in the database."""
        try:
            await mongodb.get_collection("documents").insert_one(document.dict())
            logger.info(f"Created document record: {document.document_id}")
            return document
        except Exception as e:
            logger.error(f"Error creating document record: {e}")
            raise DocumentProcessingError(f"Failed to create document record: {e}")

# Global instance
document_service = DocumentService() 